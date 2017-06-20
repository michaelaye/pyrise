import subprocess
from pathlib import Path

import click
from pptx import Presentation
from scipy.misc import imread
from six.moves.urllib.error import HTTPError
from six.moves.urllib.request import urlretrieve

from .products import PRODUCT_ID, RED_PRODUCT_ID, HiRISE_URL


def hirise_dropbox():
    return Path.home() / 'Dropbox' / 'data' / 'hirise'


def labels_root():
    return hirise_dropbox() / 'labels'


def browse_root():
    return hirise_dropbox() / 'browse'


def get_rdr_some_label(kind, obsid):
    """Download `some` PRODUCT_ID label for `obsid`.

    Note, that the RED channel is also called the B&W channel on the HiRISE website.

    Parameters
    ----------
    kind : {'RED', 'COLOR'}
        String that determines the kind of color looking for.
    obsid : str
        HiRISE obsid in the standard form of ESP_012345_1234

    Returns
    -------
    None
        Storing the label file in the `labels_root` folder.

    """
    pid = PRODUCT_ID(obsid)
    pid.kind = kind
    savepath = labels_root() / Path(pid.label_fname)
    savepath.parent.mkdir(exist_ok=True)
    print("Downloading\n", pid.label_url, 'to\n', savepath)
    try:
        urlretrieve(pid.label_url, str(savepath))
    except HTTPError as e:
        print(e)


def get_rdr_red_label(obsid):
    """Download the RED PRODUCT_ID label for `obsid`.

    Note
    ----
    The RED channel is also called the B&W channel on the HiRISE website.

    Parameters
    ----------
    obsid : str
        HiRISE obsid in the standard form of ESP_012345_1234

    Returns
    -------
    None
        Storing the label file in the `labels_root` folder.
    """
    get_rdr_some_label('RED', obsid)


def get_rdr_color_label(obsid):
    """Download the COLOR PRODUCT_ID label for `obsid`.

    Parameters
    ----------
    obsid : str
        HiRISE obsid in the standard form of ESP_012345_1234

    Returns
    -------
    None
        Storing the label file in the `labels_root` folder.
    """
    get_rdr_some_label('COLOR', obsid)


def download_product(prodid_path, saveroot=None):
    if saveroot is None:
        saveroot = hirise_dropbox()
    else:
        saveroot = Path(saveroot)
        if not saveroot.is_absolute():
            saveroot = hirise_dropbox() / saveroot

    saveroot.mkdir(exist_ok=True)
    url = HiRISE_URL(prodid_path)
    savepath = saveroot / prodid_path.name
    print("Downloading\n", url.url, 'to\n', savepath)
    try:
        urlretrieve(url.url, str(savepath))
    except HTTPError as e:
        print(e)
    return savepath


def download_RED_product(obsid, ccdno, channel, saveroot=None, overwrite=False):
    if saveroot is None:
        saveroot = hirise_dropbox()
    else:
        saveroot = Path(saveroot)
        if not saveroot.is_absolute():
            saveroot = hirise_dropbox() / saveroot

    pid = RED_PRODUCT_ID(obsid, ccdno, channel)
    savepath = saveroot / obsid / pid.fname
    savepath.parent.mkdir(parents=True, exist_ok=True)

    # if file already is there:
    if savepath.exists() and not overwrite:
        return savepath

    print("Downloading\n", pid.furl, '\nto\n', savepath)
    try:
        urlretrieve(pid.furl, str(savepath))
    except HTTPError as e:
        print(e)
    return savepath


def download_browse_product(obsid, kind='RED', annotated=True, saveroot='.', overwrite=False):
    """Download a browse product from HiRISE website.

    By default, store it in current path.

    Parameters
    ----------
    obsid : str
        HiRISE observation ID
    kind : {'RED','COLOR'}, optional
        String indicating the kind of required browse product. Default: 'RED'
    annotated : bool, optional
        Switch to control if the annotated version is required. Default: True
    saveroot : str, pathlib.Path, optional
        Path in where the browse product is being stored. Default: '.'
    overwrite : bool, optional
        Boolean switch to control if an existing path should be overwritten. Default: False
    """
    pid = PRODUCT_ID(f"{obsid}_{kind}")
    if annotated is True:
        url = pid.abrowse_url
    else:
        url = pid.browse_url
    saveroot = Path(saveroot)
    savepath = saveroot / Path(url).name
    savepath.parent.mkdir(parents=True, exist_ok=True)

    if savepath.exists() and not overwrite:
        print("Path exists. No download required.", savepath)
        return savepath

    print("Downloading\n", url, '\nto\n', savepath)
    try:
        urlretrieve(url, str(savepath))
    except HTTPError as e:
        print(e)
    return savepath


@click.command()
@click.argument('obsid')
def get_and_display_browse_product(obsid):
    path = download_browse_product(obsid)
    subprocess.run(["open", "-a", "Preview", path])


def create_browse_presentation(obsids, savename='obsid_browse_images', **kwargs):
    savepath = str(savename) + '.pptx'
    prs = Presentation()

    pic_left = int(prs.slide_width * 0.0)
    pic_top = int(prs.slide_height * 0.0)

    blank_slide_layout = prs.slide_layouts[6]

    for obsid in obsids[0:]:
        imgpath = str(download_browse_product(obsid, **kwargs))
        slide = prs.slides.add_slide(blank_slide_layout)
        img = imread(imgpath)
        ratio = img.shape[0] / img.shape[1]
        if ratio < 1:
            pic_width = int(prs.slide_width)
            pic_height = int(pic_width * img.shape[0] / img.shape[1])
        else:
            pic_height = int(prs.slide_height)
            pic_width = int(pic_height / ratio)
        slide.shapes.add_picture(imgpath, pic_left, pic_top, pic_width, pic_height)

    prs.save(savepath)
